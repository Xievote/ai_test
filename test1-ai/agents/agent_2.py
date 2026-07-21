"""
文档转换智能体 - 基于 LangChain 框架
功能：将用户原始内容转换为 Markdown，再转为 PPT/WORD/EXCEL 格式
"""

import os
import time
import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Optional, Callable

from dotenv import load_dotenv

load_dotenv()

from langchain.tools import tool
from langchain.agents import create_agent
from langchain.chat_models import init_chat_model
from langchain.messages import HumanMessage

# ============================================================
# 日志配置（仅写文件，stdout 留给进度输出）
# ============================================================
LOG_DIR = Path(__file__).parent / "logs"
LOG_DIR.mkdir(exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
    handlers=[
        logging.FileHandler(LOG_DIR / "agent_2.log", encoding="utf-8"),
    ],
)
logger = logging.getLogger("DocAgent")

# ============================================================
# 输出目录
# ============================================================
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ============================================================
# 支持的文档类型与格式
# ============================================================
SUPPORTED_DOC_TYPES = {"报告", "汇报", "总结", "清单"}
SUPPORTED_DOC_FORMATS = {"PPT", "WORD", "EXCEL"}

# ============================================================
# 内容长度限制
# ============================================================
MAX_CONTENT_LENGTH = 15000  # 原始内容最大字符数，超出将截断并警告

# ============================================================
# LLM 单例 + 带重试的调用
# ============================================================
_llm_instance = None
LLM_MODEL_ID = "deepseek:deepseek-v4-flash"
MAX_RETRIES = 3           # LLM 调用最大重试次数
RETRY_BASE_DELAY = 1.0    # 指数退避基础延迟（秒）


def get_llm():
    """获取 LLM 模块级单例，避免重复创建连接池。"""
    global _llm_instance
    if _llm_instance is None:
        _llm_instance = init_chat_model(LLM_MODEL_ID)
        logger.info(f"LLM 实例已初始化: {LLM_MODEL_ID}")
    return _llm_instance


def invoke_llm_with_retry(prompt: str, max_retries: int = MAX_RETRIES) -> str:
    """带指数退避重试的 LLM 调用。

    Args:
        prompt: 提示词文本
        max_retries: 最大重试次数

    Returns:
        LLM 响应文本

    Raises:
        RuntimeError: 所有重试均失败时抛出
    """
    last_error = None
    for attempt in range(max_retries + 1):
        try:
            llm = get_llm()
            response = llm.invoke(prompt)
            return response.content.strip()
        except Exception as e:
            last_error = e
            if attempt < max_retries:
                delay = RETRY_BASE_DELAY * (2 ** attempt)
                logger.warning(f"LLM 调用失败 (第 {attempt + 1}/{max_retries + 1} 次)，"
                               f"{delay:.1f}s 后重试: {e}")
                time.sleep(delay)
            else:
                logger.error(f"LLM 调用失败，已达最大重试次数: {e}")
    raise RuntimeError(f"LLM 调用失败（已重试 {max_retries} 次）: {last_error}")


# ============================================================
# 进度通知
# ============================================================
_progress_callback: Callable[[str], None] | None = None


def set_progress_callback(callback: Callable[[str], None] | None):
    """设置进度回调函数，用于通知外部调用方当前处理进度。

    Args:
        callback: 接收进度消息的函数，传 None 可清除回调
    """
    global _progress_callback
    _progress_callback = callback


def _notify_progress(msg: str):
    """推送进度消息。同时记录日志。"""
    logger.info(f"[进度] {msg}")
    if _progress_callback:
        _progress_callback(msg)


# 每种文档类型的 Markdown 模板提示词
DOC_TYPE_PROMPTS = {
    "报告": """请将以下内容整理为一份结构完整的**报告**，使用 Markdown 格式输出。

报告结构要求：
- 标题（一级标题，明确报告主题）
- 摘要/概述（简要说明报告背景与目的）
- 正文（分章节展开，每章使用二级标题，包含数据、分析、结论）
- 总结与建议

格式要求：使用规范的 Markdown 语法，表格用 | 分隔，重点内容用 **加粗**。""",

    "汇报": """请将以下内容整理为一份**工作汇报**，使用 Markdown 格式输出。

汇报结构要求：
- 标题（一级标题）
- 工作概述（本阶段主要工作内容）
- 重点工作进展（分条目列出，每项包含：任务名称、完成情况、关键成果）
- 数据与指标（如涉及数据，用表格呈现）
- 存在问题与风险
- 下阶段计划

格式要求：层级清晰，使用有序/无序列表，表格用 | 分隔。""",

    "总结": """请将以下内容整理为一份**总结文档**，使用 Markdown 格式输出。

总结结构要求：
- 标题（一级标题）
- 背景/概况
- 核心成果/要点（分条列出，突出关键结论）
- 经验与反思
- 后续展望

格式要求：条理清晰，要点突出，可用列表和表格辅助说明。""",

    "清单": """请将以下内容整理为一份**清单/列表**，使用 Markdown 格式输出。

清单结构要求：
- 标题（一级标题）
- 清单说明（简要说明清单用途）
- 清单正文（使用表格或有序/无序列表，每条清晰标注序号、内容、状态等）

格式要求：优先使用表格（列：序号 | 项目 | 说明/状态），确保每项清晰可辨。""",
}


# ============================================================
# 输入验证
# ============================================================
def validate_inputs(doc_type: str, doc_format: str) -> tuple[bool, str]:
    """验证用户输入的文档类型和格式是否受支持。"""
    errors = []
    if doc_type not in SUPPORTED_DOC_TYPES:
        errors.append(f"不支持的文档类型「{doc_type}」，支持的类型：{', '.join(SUPPORTED_DOC_TYPES)}")
    if doc_format.upper() not in SUPPORTED_DOC_FORMATS:
        errors.append(f"不支持的文档格式「{doc_format}」，支持的格式：{', '.join(SUPPORTED_DOC_FORMATS)}")
    if errors:
        return False, "；".join(errors)
    return True, ""


# ============================================================
# 工具 1: 将原始内容转换为 Markdown
# ============================================================
@tool
def convert_to_markdown(raw_content: str, doc_type: str) -> str:
    """将用户原始内容转换为结构完整、格式规范的 Markdown 文档。

    Args:
        raw_content: 用户提供的原始文本内容
        doc_type: 文档类型，可选值：报告、汇报、总结、清单
    """
    _notify_progress(f"[1/2] 正在分析原始内容，文档类型: {doc_type}...")

    valid, error_msg = validate_inputs(doc_type, "PPT")  # 仅验 doc_type
    if not valid and "格式" not in error_msg:
        return f"[错误] {error_msg}"

    if not raw_content or not raw_content.strip():
        return "[错误] 原始内容不能为空，请提供有效的文本内容。"

    # 长内容截断 + 警告
    truncated = False
    if len(raw_content) > MAX_CONTENT_LENGTH:
        trunc_msg = (
            f"⚠ 原始内容过长（{len(raw_content)} 字符），"
            f"已截断至 {MAX_CONTENT_LENGTH} 字符。建议精简输入以获得更完整的文档。"
        )
        logger.warning(trunc_msg)
        _notify_progress(trunc_msg)
        raw_content = raw_content[:MAX_CONTENT_LENGTH]
        truncated = True

    try:
        prompt = DOC_TYPE_PROMPTS.get(doc_type, DOC_TYPE_PROMPTS["报告"])

        full_prompt = f"""{prompt}

用户原始内容如下：
---
{raw_content.strip()}
---

请严格按照上述要求输出完整的 Markdown 文档，不要省略任何内容。"""

        _notify_progress(f"[1/2] 正在生成 Markdown 文档（使用 {LLM_MODEL_ID}）...")
        markdown_content = invoke_llm_with_retry(full_prompt)
        result_len = len(markdown_content)
        logger.info(f"Markdown 生成成功，长度: {result_len} 字符")
        _notify_progress(f"[1/2] ✓ Markdown 生成完成 ({result_len} 字符)")

        if truncated:
            markdown_content = (
                f"> ⚠ **注意**：原始内容过长已被截断，本 Markdown 仅基于前 "
                f"{MAX_CONTENT_LENGTH} 字符生成。\n\n{markdown_content}"
            )
        return markdown_content

    except RuntimeError as e:
        logger.error(f"Markdown 生成失败: {e}")
        return f"[错误] Markdown 生成失败: {str(e)}"


# ============================================================
# 工具 2: Markdown → PPT
# ============================================================

# ---- 每页内容容量常量 ----
MAX_BULLETS_PER_SLIDE = 7      # 每页最多容纳的列表项数
MAX_TEXT_LINES_PER_SLIDE = 10   # 每页最多容纳的文本行数（含标题）
MAX_TABLE_ROWS_PER_SLIDE = 8    # 每页表格最多行数


def _estimate_text_height(lines_count: int, heading_count: int = 0, font_pt: int = 16) -> float:
    """估算文本区域的高度（英寸），用于判断是否超出页面范围。"""
    line_spacing = font_pt * 1.5 / 72   # 每行高度（英寸）
    heading_spacing = heading_count * 28 / 72  # 标题额外间距
    return lines_count * line_spacing + heading_spacing


def _generate_title_from_markdown(markdown_content: str) -> str:
    """调用大模型分析文档内容，生成不超过15个汉字的简洁标题作为文件名。

    Args:
        markdown_content: 完整的 Markdown 文本内容

    Returns:
        清理后的标题字符串（≤15个汉字）
    """
    _notify_progress("[标题] 正在分析文档内容生成文件名...")
    try:
        prompt = f"""请根据以下文档内容，总结出一个简洁的标题，要求：
1. 标题不超过 15 个汉字
2. 概括文档的核心主题
3. 只输出标题本身，不要带任何引号、标点或额外说明

文档内容：
{markdown_content[:3000]}"""

        title = invoke_llm_with_retry(prompt, max_retries=2)

        # 清理：去除引号、多余标点、换行
        title = re.sub(r'[\"\'""''「」『』\n\r]', '', title)
        title = re.sub(r'[\.。！!？?，,、：:；;…—\-\(\)（）\[\]【】]', '', title)
        title = title.strip()

        # 限制15个汉字
        if len(title) > 15:
            title = title[:15]

        # 如果标题为空或无效，使用默认标题
        if not title or len(title.strip()) < 2:
            title = "文档"

        logger.info(f"LLM 生成的标题: {title}")
        _notify_progress(f"[标题] ✓ 文档标题: {title}")
        return title
    except RuntimeError as e:
        logger.warning(f"标题生成失败，使用默认标题: {e}")
        _notify_progress(f"[标题] ⚠ 标题生成失败，使用默认文件名")
        return "文档"


def _parse_markdown_to_slides(markdown_content: str) -> list[dict]:
    """将 Markdown 智能解析为多页幻灯片结构。

    分页策略：
    - 标题页：提取第一个 H1 作为封面
    - 目录页：汇总所有 H2 标题作为目录
    - 内容页：每个 H2 下级内容为一页；根据容量阈值自动拆分
    - 表格页：表格数据按行数阈值拆分
    - 结束页：Q&A/感谢观看
    """
    slides = []
    lines = markdown_content.strip().split("\n")

    # ---------- 第一阶段：解析为段落结构 ----------
    # 结构: {"sections": [
    #   {"type": "h1", "text": "...", "children": [...]},
    #   {"type": "h2", "text": "...", "children": [...]},
    #   ...
    # ]}
    sections = []
    current_section = None

    def _flush_section():
        nonlocal current_section
        if current_section is not None:
            sections.append(current_section)
            current_section = None

    for line in lines:
        # 一级标题
        if line.startswith("# ") and not line.startswith("## "):
            _flush_section()
            current_section = {
                "type": "h1",
                "text": line.lstrip("# ").strip(),
                "items": [],
            }
        # 二级标题
        elif line.startswith("## ") and not line.startswith("### "):
            _flush_section()
            current_section = {
                "type": "h2",
                "text": re.sub(r"^#{2,3}\s", "", line).strip(),
                "items": [],
            }
        # 三级标题
        elif line.startswith("### "):
            _flush_section()
            current_section = {
                "type": "h3",
                "text": re.sub(r"^#{2,3}\s", "", line).strip(),
                "items": [],
            }
        # 表格行
        elif line.strip().startswith("|") and line.strip().endswith("|"):
            if "---" not in line:
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                if current_section is None:
                    current_section = {"type": "body", "text": "", "items": []}
                current_section["items"].append(("table_row", cells))
        # 列表项
        elif re.match(r"^[\s]*[-*+]\s", line) or re.match(r"^\d+[\.)]\s", line):
            item_text = re.sub(r"^[\s]*[-*+\d]+[\.)]\s*", "", line).strip()
            if item_text:
                if current_section is None:
                    current_section = {"type": "body", "text": "", "items": []}
                current_section["items"].append(("bullet", item_text))
        # 空行
        elif not line.strip():
            continue
        # 普通文本
        else:
            clean = re.sub(r"(\*\*|__)(.*?)\1", r"\2", line.strip())
            if clean:
                if current_section is None:
                    current_section = {"type": "body", "text": "", "items": []}
                current_section["items"].append(("text", clean))

    _flush_section()

    if not sections:
        return [{"page_type": "content", "title": "文档", "items": [("text", markdown_content)]}]

    # ---------- 第二阶段：构建幻灯片 ----------
    h1_title = ""
    h2_titles = []

    for sec in sections:
        if sec["type"] == "h1":
            h1_title = sec["text"]
        elif sec["type"] in ("h2", "h3"):
            h2_titles.append(sec["text"])

    # 1) 封面页
    if h1_title:
        slides.append({
            "page_type": "cover",
            "title": h1_title,
            "items": [],
        })

    # 2) 目录页（如果有 ≥2 个二级标题）
    if len(h2_titles) >= 2:
        slides.append({
            "page_type": "toc",
            "title": "目录",
            "toc_items": h2_titles,
        })

    # 3) 内容页：按 H2/H3 分段，每段内容做容量检测与拆分
    for sec in sections:
        if sec["type"] == "h1":
            continue

        sec_title = sec["text"]
        sec_items = sec["items"]

        if not sec_items:
            continue

        # 分离表格和非表格内容
        table_items = [it for it in sec_items if it[0] == "table_row"]
        non_table = [it for it in sec_items if it[0] != "table_row"]

        # 3a) 非表格内容：检测高度并拆分
        if non_table:
            sub_slides = _split_items_by_capacity(non_table, sec_title)
            slides.extend(sub_slides)

        # 3b) 表格内容：按行数阈值拆分
        if table_items:
            table_slides = _split_table_items(table_items, sec_title)
            slides.extend(table_slides)

    # 4) 结束页
    slides.append({
        "page_type": "ending",
        "title": "感谢观看",
        "items": [],
    })

    return slides


def _split_items_by_capacity(items: list, section_title: str) -> list[dict]:
    """将内容项按容量阈值拆分为多页幻灯片。"""
    result = []
    chunk = []
    heading_in_chunk = False
    line_count = 0

    def _emit_chunk():
        nonlocal chunk, heading_in_chunk, line_count
        if chunk:
            title = section_title
            result.append({
                "page_type": "content",
                "title": title,
                "items": list(chunk),
            })
            chunk = []
            heading_in_chunk = False
            line_count = 0

    for item in items:
        t, val = item
        if t == "heading":
            # heading 出现在中间 → 先输出当前 chunk，然后开新页
            _emit_chunk()
            result.append({
                "page_type": "content",
                "title": val,
                "items": [],
            })
        elif t == "bullet":
            chunk.append(item)
            line_count += 1
            if line_count >= MAX_BULLETS_PER_SLIDE:
                _emit_chunk()
        elif t == "text":
            chunk.append(item)
            line_count += max(1, len(val) // 60)  # 长文本行数估算
            if line_count >= MAX_TEXT_LINES_PER_SLIDE:
                _emit_chunk()

    _emit_chunk()
    return result


def _split_table_items(table_items: list, section_title: str) -> list[dict]:
    """将表格行按 MAX_TABLE_ROWS_PER_SLIDE 拆分为多页。"""
    result = []
    for i in range(0, len(table_items), MAX_TABLE_ROWS_PER_SLIDE):
        chunk = table_items[i:i + MAX_TABLE_ROWS_PER_SLIDE]
        result.append({
            "page_type": "content",
            "title": section_title if i == 0 else f"{section_title}（续）",
            "items": list(chunk),
        })
    return result


def _render_slide(prs, slide_data: dict):
    """在 Presentation 对象中渲染一页幻灯片。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    page_type = slide_data.get("page_type", "content")
    title = slide_data.get("title", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # ---- 封面页 ----
    if page_type == "cover":
        # 主标题居中
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "文档转换智能体 · 自动生成"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p2.alignment = PP_ALIGN.CENTER

        # 底部分隔线
        line_shape = slide.shapes.add_shape(
            1, Inches(3.0), Inches(5.8), Inches(7.3), Pt(2)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        line_shape.line.fill.background()
        return

    # ---- 目录页 ----
    if page_type == "toc":
        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        p.alignment = PP_ALIGN.LEFT

        # 分隔线
        slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.7), Pt(3)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        slide.shapes[-1].line.fill.background()

        # 目录项
        toc_items = slide_data.get("toc_items", [])
        if toc_items:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(5.0))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for idx, item_text in enumerate(toc_items):
                if idx == 0:
                    p = tf2.paragraphs[0]
                else:
                    p = tf2.add_paragraph()
                p.text = f"{idx + 1}.  {item_text}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_before = Pt(14)
                p.space_after = Pt(6)
        return

    # ---- 结束页 ----
    if page_type == "ending":
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "Q & A"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p2.alignment = PP_ALIGN.CENTER
        return

    # ---- 内容页 ----
    # 标题区域
    if title:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        p.alignment = PP_ALIGN.LEFT

        # 分隔线
        slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(11.7), Pt(3)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
        slide.shapes[-1].line.fill.background()

    # 内容区域
    content_top = Inches(1.8) if title else Inches(0.4)
    content_left = Inches(1.0)
    content_width = Inches(11.3)
    content_height = Inches(5.2) if title else Inches(6.6)

    txBox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    items = slide_data.get("items", [])
    table_rows = [it[1] for it in items if it[0] == "table_row"]
    non_table = [it for it in items if it[0] != "table_row"]

    # 渲染非表格内容
    first = True
    for item_type, item_value in non_table:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if item_type == "heading":
            p.text = item_value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_before = Pt(14)
            p.space_after = Pt(6)
        elif item_type == "bullet":
            p.text = f"• {item_value}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            p.space_before = Pt(5)
            p.space_after = Pt(2)
        elif item_type == "text":
            p.text = item_value
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            p.space_before = Pt(6)
            p.space_after = Pt(3)

    # 清除空的首段落
    if len(tf.paragraphs) > 1:
        first_p = tf.paragraphs[0]
        if not first_p.text.strip() and first_p._p is not None:
            try:
                first_p._p.getparent().remove(first_p._p)
            except Exception:
                pass

    # 渲染表格
    if table_rows:
        num_cols = max(len(r) for r in table_rows) if table_rows else 1
        num_rows = len(table_rows)
        table_top = Inches(3.8) if non_table else Inches(1.8)
        table_height = min(num_rows * 0.45, 3.2)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(1.0), Inches(table_top),
            Inches(11.3), Inches(table_height)
        )
        table = table_shape.table

        for r_idx, row in enumerate(table_rows):
            for c_idx, cell_text in enumerate(row):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        if r_idx == 0:
                            paragraph.font.bold = True


@tool
def convert_markdown_to_ppt(markdown_content: str) -> str:
    """将 Markdown 内容转换为多页 PPT 文件，自动智能分页并生成文件名。

    Args:
        markdown_content: 完整的 Markdown 文本内容
    """
    _notify_progress("[PPT] 开始转换为 PPT 格式（智能分页 + 标题生成）...")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        # 1. 调用 LLM 生成文档标题（用作文件名）
        doc_title = _generate_title_from_markdown(markdown_content)
        if not doc_title or len(doc_title.strip()) < 2:
            doc_title = "文档"

        # 2. 解析 Markdown → 多页幻灯片结构
        _notify_progress("[PPT] 正在解析文档结构并智能分页...")
        slides_data = _parse_markdown_to_slides(markdown_content)
        total_slides = len(slides_data)
        logger.info(f"幻灯片解析完成，共 {total_slides} 页")
        _notify_progress(f"[PPT] 解析完成，共 {total_slides} 页，正在渲染...")

        # 3. 创建 PPT 并渲染每一页
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for idx, slide_data in enumerate(slides_data, 1):
            _notify_progress(f"[PPT] 渲染第 {idx}/{total_slides} 页...")
            _render_slide(prs, slide_data)

        # 4. 保存文件（使用 LLM 生成的标题作为文件名）
        safe_title = re.sub(r'[\\/:*?"<>|]', '', doc_title)
        safe_title = safe_title.replace(' ', '_') or "文档"
        filename = f"{safe_title}.pptx"
        filepath = OUTPUT_DIR / filename

        # 文件名冲突时添加序号
        counter = 1
        while filepath.exists():
            stem = safe_title
            filename = f"{stem}_{counter}.pptx"
            filepath = OUTPUT_DIR / filename
            counter += 1

        _notify_progress("[PPT] 正在保存文件...")
        prs.save(str(filepath))

        logger.info(f"PPT 文件已保存 ({total_slides} 页): {filepath}")
        _notify_progress(f"[PPT] ✓ 文件已保存: {filename}")
        return (
            f"[成功] PPT 文件已生成\n"
            f"  • 页数：{total_slides} 页\n"
            f"  • 文件名：{filename}\n"
            f"  • 保存路径：{filepath}"
        )

    except Exception as e:
        logger.error(f"PPT 转换失败: {e}")
        return f"[错误] PPT 文件生成失败: {str(e)}"


# ============================================================
# 工具 3: Markdown → WORD
# ============================================================
@tool
def convert_markdown_to_word(markdown_content: str) -> str:
    """将 Markdown 内容转换为 WORD 文件。

    Args:
        markdown_content: 完整的 Markdown 文本内容
    """
    _notify_progress("[WORD] 开始转换为 WORD 格式...")

    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT

        doc = Document()

        # 设置默认字体
        style = doc.styles["Normal"]
        font = style.font
        font.name = "微软雅黑"
        font.size = Pt(11)

        lines = markdown_content.strip().split("\n")
        i = 0

        while i < len(lines):
            line = lines[i]

            # 一级标题
            if line.startswith("# ") and not line.startswith("## "):
                text = line.lstrip("# ").strip()
                doc.add_heading(text, level=1)
                i += 1
                continue

            # 二级标题
            if line.startswith("## ") and not line.startswith("### "):
                text = line.lstrip("#").strip()
                doc.add_heading(text, level=2)
                i += 1
                continue

            # 三级标题
            if line.startswith("### "):
                text = line.lstrip("#").strip()
                doc.add_heading(text, level=3)
                i += 1
                continue

            # 表格
            if line.strip().startswith("|") and line.strip().endswith("|"):
                table_data = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    row_line = lines[i].strip()
                    if "---" not in row_line.replace("-", ""):
                        cells = [c.strip() for c in row_line.strip("|").split("|")]
                        table_data.append(cells)
                    i += 1

                if table_data:
                    num_cols = max(len(r) for r in table_data)
                    num_rows = len(table_data)
                    table = doc.add_table(rows=num_rows, cols=num_cols, style="Light Grid Accent 1")
                    table.alignment = WD_TABLE_ALIGNMENT.CENTER

                    for r_idx, row in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row):
                            if c_idx < num_cols:
                                cell = table.cell(r_idx, c_idx)
                                cell.text = cell_text
                                for paragraph in cell.paragraphs:
                                    paragraph.paragraph_format.space_before = Pt(2)
                                    paragraph.paragraph_format.space_after = Pt(2)
                                    for run in paragraph.runs:
                                        run.font.size = Pt(10)
                                        if r_idx == 0:
                                            run.font.bold = True
                    doc.add_paragraph()  # 表格后空行
                continue

            # 无序列表
            if re.match(r"^[\s]*[-*+]\s", line):
                text = re.sub(r"^[\s]*[-*+]\s*", "", line)
                p = doc.add_paragraph(text, style="List Bullet")
                for run in p.runs:
                    run.text = re.sub(r"\*\*(.+?)\*\*", r"\1", run.text)
                i += 1
                continue

            # 有序列表
            if re.match(r"^\d+[\.)]\s", line):
                text = re.sub(r"^\d+[\.)]\s*", "", line)
                p = doc.add_paragraph(text, style="List Number")
                for run in p.runs:
                    run.text = re.sub(r"\*\*(.+?)\*\*", r"\1", run.text)
                i += 1
                continue

            # 空行
            if not line.strip():
                i += 1
                continue

            # 普通段落
            clean_text = line.strip()
            p = doc.add_paragraph(clean_text)
            for run in p.runs:
                run.text = re.sub(r"\*\*(.+?)\*\*", r"\1", run.text)
            i += 1

        # 使用 LLM 生成的标题作为文件名
        doc_title = _generate_title_from_markdown(markdown_content)
        safe_title = re.sub(r'[\\/:*?"<>|]', '', doc_title)
        safe_title = safe_title.replace(' ', '_') or "文档"
        filename = f"{safe_title}.docx"
        filepath = OUTPUT_DIR / filename

        counter = 1
        while filepath.exists():
            filename = f"{safe_title}_{counter}.docx"
            filepath = OUTPUT_DIR / filename
            counter += 1

        _notify_progress("[WORD] 正在保存文件...")
        doc.save(str(filepath))

        logger.info(f"WORD 文件已保存: {filepath}")
        _notify_progress(f"[WORD] ✓ 文件已保存: {filename}")
        return (
            f"[成功] WORD 文件已生成\n"
            f"  • 文件名：{filename}\n"
            f"  • 保存路径：{filepath}"
        )

    except Exception as e:
        logger.error(f"WORD 转换失败: {e}")
        return f"[错误] WORD 文件生成失败: {str(e)}"


# ============================================================
# 工具 4: Markdown → EXCEL
# ============================================================
@tool
def convert_markdown_to_excel(markdown_content: str) -> str:
    """将 Markdown 内容转换为 EXCEL 文件。

    Args:
        markdown_content: 完整的 Markdown 文本内容
    """
    _notify_progress("[EXCEL] 开始转换为 EXCEL 格式...")

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        header_font = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1A3C6E", end_color="1A3C6E", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell_font = Font(name="微软雅黑", size=10)
        cell_alignment = Alignment(vertical="top", wrap_text=True)
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        lines = markdown_content.strip().split("\n")
        current_row = 1
        i = 0

        while i < len(lines):
            line = lines[i]

            # 一级标题 → 作为 Sheet 标题行
            if line.startswith("# ") and not line.startswith("## "):
                text = line.lstrip("# ").strip()
                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=8)
                cell = ws.cell(row=current_row, column=1, value=text)
                cell.font = Font(name="微软雅黑", size=16, bold=True, color="1A3C6E")
                cell.alignment = Alignment(horizontal="center", vertical="center")
                current_row += 2
                i += 1
                continue

            # 二/三级标题
            if line.startswith("## ") or line.startswith("### "):
                text = re.sub(r"^#{2,3}\s", "", line).strip()
                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=8)
                cell = ws.cell(row=current_row, column=1, value=text)
                cell.font = Font(name="微软雅黑", size=12, bold=True, color="333333")
                cell.alignment = Alignment(horizontal="left", vertical="center")
                current_row += 1
                i += 1
                continue

            # 表格
            if line.strip().startswith("|") and line.strip().endswith("|"):
                table_data = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    row_line = lines[i].strip()
                    if "---" not in row_line.replace("-", ""):
                        cells = [c.strip() for c in row_line.strip("|").split("|")]
                        table_data.append(cells)
                    i += 1

                if table_data:
                    for r_idx, row in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row):
                            excel_cell = ws.cell(row=current_row + r_idx, column=c_idx + 1, value=cell_text)
                            excel_cell.font = cell_font
                            excel_cell.alignment = cell_alignment
                            excel_cell.border = thin_border
                            if r_idx == 0:
                                excel_cell.font = header_font
                                excel_cell.fill = header_fill
                                excel_cell.alignment = header_alignment

                    # 自动调整列宽
                    for c_idx in range(len(table_data[0]) if table_data else 0):
                        max_width = max(
                            len(str(row[c_idx])) if c_idx < len(row) else 0
                            for row in table_data
                        )
                        ws.column_dimensions[get_column_letter(c_idx + 1)].width = min(max_width * 2 + 4, 40)

                    current_row += len(table_data) + 1
                continue

            # 列表项
            if re.match(r"^[\s]*[-*+]\s", line) or re.match(r"^\d+[\.)]\s", line):
                text = re.sub(r"^[\s]*[-*+\d]+[\.)]\s*", "", line).strip()
                ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=8)
                cell = ws.cell(row=current_row, column=1, value=f"• {text}")
                cell.font = cell_font
                cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
                current_row += 1
                i += 1
                continue

            # 空行
            if not line.strip():
                current_row += 1
                i += 1
                continue

            # 普通文本
            ws.merge_cells(start_row=current_row, start_column=1, end_row=current_row, end_column=8)
            cell = ws.cell(row=current_row, column=1, value=line.strip())
            cell.font = cell_font
            cell.alignment = cell_alignment
            current_row += 1
            i += 1

        # 设置列宽默认值
        for col_idx in range(1, 9):
            ws.column_dimensions[get_column_letter(col_idx)].width = max(
                ws.column_dimensions[get_column_letter(col_idx)].width or 14, 14
            )

        # 使用 LLM 生成的标题作为文件名
        doc_title = _generate_title_from_markdown(markdown_content)
        safe_title = re.sub(r'[\\/:*?"<>|]', '', doc_title)
        safe_title = safe_title.replace(' ', '_') or "文档"
        filename = f"{safe_title}.xlsx"
        filepath = OUTPUT_DIR / filename

        counter = 1
        while filepath.exists():
            filename = f"{safe_title}_{counter}.xlsx"
            filepath = OUTPUT_DIR / filename
            counter += 1

        _notify_progress("[EXCEL] 正在保存文件...")
        wb.save(str(filepath))

        logger.info(f"EXCEL 文件已保存: {filepath}")
        _notify_progress(f"[EXCEL] ✓ 文件已保存: {filename}")
        return (
            f"[成功] EXCEL 文件已生成\n"
            f"  • 文件名：{filename}\n"
            f"  • 保存路径：{filepath}"
        )

    except Exception as e:
        logger.error(f"EXCEL 转换失败: {e}")
        return f"[错误] EXCEL 文件生成失败: {str(e)}"


# ============================================================
# 创建 LangChain Agent
# ============================================================
SYSTEM_PROMPT = """你是一个专业的文档转换智能体。你的任务是帮助用户将原始内容转换为指定格式的文档。

工作流程（必须严格按顺序执行）：
1. 首先确认用户提供的三要素：原始内容、文档类型（报告/汇报/总结/清单）、目标格式（PPT/WORD/EXCEL）
2. 调用 convert_to_markdown 工具，将原始内容转换为结构化的 Markdown 文档
3. 根据目标格式，调用对应的转换工具：
   - PPT  → convert_markdown_to_ppt
   - WORD → convert_markdown_to_word
   - EXCEL → convert_markdown_to_excel
4. 将工具返回的结果（文件保存路径或错误信息）告知用户

重要规则：
- 必须先转换为 Markdown，再转为目标格式
- 如果用户没有提供文档类型，默认使用"报告"
- 如果用户没有提供目标格式，询问用户后再继续
- 如果任何步骤出现 [错误] 前缀的返回，停止后续步骤并告知用户错误原因
- 工具调用完成且成功时，直接告诉用户结果，不要再次调用工具
"""

model = get_llm()

agent = create_agent(
    model=model,
    tools=[
        convert_to_markdown,
        convert_markdown_to_ppt,
        convert_markdown_to_word,
        convert_markdown_to_excel,
    ],
    system_prompt=SYSTEM_PROMPT,
)

# ============================================================
# 用户交互接口
# ============================================================
def run_conversion(
    raw_content: str,
    doc_type: str = "报告",
    doc_format: str = "WORD",
    verbose: bool = True,
) -> str:
    """执行文档转换的入口函数。

    Args:
        raw_content: 用户原始文本内容
        doc_type: 文档类型（报告/汇报/总结/清单）
        doc_format: 目标格式（PPT/WORD/EXCEL）
        verbose: 是否启用进度输出（打印到 stdout）

    Returns:
        Agent 执行结果字符串
    """
    # 设置进度回调（verbose 模式下打印到 stdout）
    if verbose:
        set_progress_callback(lambda msg: print(msg))
    else:
        set_progress_callback(None)

    _notify_progress(f"收到转换请求 - 类型: {doc_type}, 格式: {doc_format}")

    # 输入验证
    valid, error_msg = validate_inputs(doc_type, doc_format.upper())
    if not valid:
        logger.error(f"输入验证失败: {error_msg}")
        _notify_progress(f"✗ {error_msg}")
        return f"[错误] {error_msg}"

    doc_format = doc_format.upper()

    # 构建用户消息
    user_message = f"""请帮我将以下内容转换为文档：

【原始内容】
{raw_content.strip()}

【文档类型】{doc_type}
【目标格式】{doc_format}

请按照工作流程依次处理：先转为 Markdown，再转为 {doc_format} 文件。"""

    _notify_progress("开始执行文档转换流程...")
    try:
        inputs = {"messages": [HumanMessage(content=user_message)]}
        result = agent.invoke(inputs)
        final_msg = result["messages"][-1].content
        logger.info(f"转换任务完成: {final_msg[:100]}...")
        _notify_progress("文档转换流程结束")
        return final_msg
    except Exception as e:
        logger.error(f"Agent 执行异常: {e}")
        _notify_progress(f"✗ 文档转换过程发生异常: {e}")
        return f"[错误] 文档转换过程发生异常: {str(e)}"


def interactive_mode():
    """交互式命令行模式，引导用户逐步输入。"""
    print("=" * 60)
    print("     文档转换智能体（基于 LangChain）")
    print("=" * 60)
    print(f"支持的文档类型: {', '.join(SUPPORTED_DOC_TYPES)}")
    print(f"支持的输出格式: {', '.join(SUPPORTED_DOC_FORMATS)}")
    print(f"输出文件目录: {OUTPUT_DIR}")
    print("=" * 60)
    print()

    # 1. 获取原始内容
    print("【步骤 1/3】请输入原始内容（输入 END 结束多行输入）：")
    lines = []
    while True:
        line = input()
        if line.strip().upper() == "END":
            break
        lines.append(line)
    raw_content = "\n".join(lines)

    if not raw_content.strip():
        print("[错误] 原始内容不能为空！")
        return

    # 2. 获取文档类型
    print(f"\n【步骤 2/3】请选择文档类型：{' / '.join(SUPPORTED_DOC_TYPES)}（默认: 报告）")
    doc_type = input("> ").strip()
    if not doc_type:
        doc_type = "报告"

    # 3. 获取目标格式
    print(f"\n【步骤 3/3】请选择输出格式：{' / '.join(SUPPORTED_DOC_FORMATS)}")
    doc_format = input("> ").strip().upper()

    if doc_format not in SUPPORTED_DOC_FORMATS:
        print(f"[错误] 不支持的格式「{doc_format}」，支持：{', '.join(SUPPORTED_DOC_FORMATS)}")
        return

    print("\n正在处理，请稍候...\n")
    result = run_conversion(raw_content, doc_type, doc_format)
    print("\n" + "=" * 60)
    print(result)
    print("=" * 60)


# ============================================================
# 程序入口
# ============================================================
if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1 and sys.argv[1] == "--demo":
        # 演示模式：使用内置示例快速测试
        logger.info("启动演示模式")
        demo_content = """2024年Q3季度，我们部门完成了以下主要工作：
1. 产品A的研发已进入测试阶段，完成了5个核心模块的开发，代码覆盖率达到85%
2. 市场推广方面，在3个重点城市开展了线下推广活动，新增用户2.3万人
3. 团队建设方面，新增了2名高级工程师，组织了4次技术分享会

存在的问题：
- 产品B的进度比计划延迟了2周，主要原因是需求变更频繁
- 服务器成本相比上季度上涨了15%，需要优化资源配置

下季度计划：
- 产品A正式发布上线
- 启动产品C的预研工作
- 优化服务器架构，目标降低成本10%"""

        print("=" * 60)
        print("     演示模式 - 依次测试三种格式")
        print("=" * 60)

        for fmt in ["PPT", "WORD", "EXCEL"]:
            print(f"\n{'─' * 40}")
            print(f">>> 测试格式: {fmt}")
            print(f"{'─' * 40}")
            result = run_conversion(demo_content, "汇报", fmt)
            print(result)
    else:
        interactive_mode()
